from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class Poweregen:
  def __init__(self):
    """
    Generate a class
    Adding standard metrics (slide size)
    """
    self.prs = Presentation()
    #Set the size to 16:9
    self.slide_width = Cm(25.4)
    self.slide_height = Cm(14.29)
    self.prs.slide_width = self.slide_width
    self.prs.slide_height = self.slide_height

  def add_title(self, title, subtitle=None):
    """Create the Title page for Power Point"""
    prs = self.prs
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title # Main title
    slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(36)
    tle = slide.placeholders[0] #the placeholder for title
    #Layout adjustment to 16:9
    tle.width = Cm(18.4)
    tle.height = Cm(5)
    tle.top = Cm(3)
    tle.left = Cm(3.5)
    if subtitle != None:      # Make sure none object is not assigned to .text
      sub = slide.placeholders[1]
      #Layout adjustment to 16:9
      sub.text = subtitle # Subtitle
      sub.text_frame.paragraphs[0].font.size = Pt(24) # Set the subtitle font smaller
      sub.width = Cm(18.4)
      sub.height = Cm(3.54)
      sub.top = Cm(8)
      sub.left = Cm(3.5)

  def add_content_bullet(self, title, content, double_spacing=False, font_size=24):
    """
    This is the part adding bulletpoint style slide
    Add feature: Double-Spacing
    Add feature: Font-size control
    """
    prs = self.prs
    slide_layout = prs.slide_layouts[1]  # Assuming this is the correct layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(36) #Setup the title size

    #Layout adjustment to 16:9
    tle = slide.placeholders[0]
    tle.width = Cm(22.86)
    tle.height = Cm(2.15)
    tle.top = Cm(0.6)
    tle.left = Cm(1.27)

    #Layout of content field adjust to 16:9
    cont = slide.placeholders[1]
    cont.width = Cm(22.86)
    cont.height = Cm(10.75)
    cont.top = Cm(3.1)
    cont.left = Cm(1.27)

    tf = slide.placeholders[1].text_frame

    sep_content = content.split('\n')  # Splitting the content into separate bullet points
    tf.text = sep_content[0]
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.LEFT
    p.level = 0
    if double_spacing:
      tf.add_paragraph()

    for line in sep_content[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT
        p.level = 0

        if double_spacing:
            tf.add_paragraph()

  def add_content_subbullet(self, title, bullets, subbullets, double_spacing=False, font_size=24):
    """
    title: The big title on top
    bullets: You would need a list of big bullets you want to list
    subbullets: Then a list of list, each list within the biglist would be the subbullet points map to content in bullets
    normal stylization: Double_spacing, font_size
    """
    prs = self.prs
    slide_layout = prs.slide_layouts[1]  # Assuming this is the layout you want
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(36) #Setup the title size

    #Layout adjustment to 16:9
    tle = slide.placeholders[0]
    tle.width = Cm(22.86)
    tle.height = Cm(2.15)
    tle.top = Cm(0.6)
    tle.left = Cm(1.27)

    #Layout of content field adjust to 16:9
    cont = slide.placeholders[1]
    cont.width = Cm(22.86)
    cont.height = Cm(10.75)
    cont.top = Cm(3.1)
    cont.left = Cm(1.27)

    tf = slide.placeholders[1].text_frame  # Assuming this is the correct placeholder

    # Iterate through each bullet and its corresponding sub-bullets
    for i, bullet in enumerate(bullets):
        # Add main bullet
        if i != 0:
          p = tf.add_paragraph()
          p.text = bullet
          p.font.size = Pt(font_size)
          p.alignment = PP_ALIGN.LEFT
          p.level = 0  # Main bullet level
        else:
          tf.text = bullet
          p = tf.paragraphs[0].runs[0]
          p.font.size = Pt(font_size)
          p.alignment = PP_ALIGN.LEFT
          p.level = 0  # Main bullet level

        # Add sub-bullets
        if i < len(subbullets):
          for subbullet in subbullets[i]:
              p = tf.add_paragraph()
              p.text = subbullet
              p.font.size = Pt(font_size)
              p.alignment = PP_ALIGN.LEFT
              p.level = 1  # Sub-bullet level

        if double_spacing:
            tf.add_paragraph()  # Add an extra paragraph for double spacing

  def add_comparison_slide(self, title, left_title, left_content, right_title, right_content, font_size=24):
    #Font setup util
    def set_paragraph_font_size(paragraph, font_sizee):
        for run in paragraph.runs:
            run.font.size = Pt(font_sizee)
    """
    Add a comparison slide to the presentation.
    title: Title of the slide
    left_title: Subtitle for the left side
    right_title: Subtitle for the right side
    left_content: Content for the left side, separated by new lines for bullet points
    right_content: Content for the right side, separated by new lines for bullet points
    font_size: Font size of the content
    """
    prs = self.prs
    slide_layout = prs.slide_layouts[4]  # Assuming this is the comparison layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[0].text_frame.paragraphs[0].font.size = Pt(36) #Setup the title size

    #Layout adjustment to 16:9
    tle = slide.placeholders[0]
    tle.width = Cm(22.86)
    tle.height = Cm(2.15)
    tle.top = Cm(0.6)
    tle.left = Cm(1.27)

    # Assuming the first two placeholders are for the comparison text
    left_titlebox = slide.placeholders[1]
    left_textbox = slide.placeholders[2]
    right_titlebox = slide.placeholders[3]
    right_textbox = slide.placeholders[4]


    #Layout adjustment to 16:9 for left right title boxes
    left_titlebox.width, right_titlebox.width = Cm(11.22), Cm(11.22)
    left_titlebox.height, right_titlebox.height = Cm(1.18), Cm(1.18)
    left_titlebox.left, right_titlebox.left = Cm(1.27), Cm(12.91)
    left_titlebox.top, right_titlebox.top = Cm(3.1), Cm(3.1)
    #Layout adjustment to 16:9 for left right text boxes
    left_textbox.width, right_textbox.width = Cm(11.22), Cm(11.22)
    left_textbox.height, right_textbox.height = Cm(8.98), Cm(8.98)
    left_textbox.left, right_textbox.left = Cm(1.27), Cm(12.91)
    left_textbox.top, right_textbox.top = Cm(4.6), Cm(4.6)


    # Set font size for title boxes
    for title_box in [left_titlebox, right_titlebox]:
        title_box.text = left_title if title_box == left_titlebox else right_title
        for paragraph in title_box.text_frame.paragraphs:
            set_paragraph_font_size(paragraph, font_size)

    # Add left side content and set font size
    left_tf = left_textbox.text_frame
    left_tf.text = left_content.split('\n')[0]
    set_paragraph_font_size(left_tf.paragraphs[0], font_size)

    for line in left_content.split('\n')[1:]:
        p = left_tf.add_paragraph()
        p.text = line
        set_paragraph_font_size(p, font_size)

    # Add right side content and set font size
    right_tf = right_textbox.text_frame
    right_tf.text = right_content.split('\n')[0]
    set_paragraph_font_size(right_tf.paragraphs[0], font_size)

    for line in right_content.split('\n')[1:]:
        p = right_tf.add_paragraph()
        p.text = line
        set_paragraph_font_size(p, font_size)

  def add_picture_slide(self, slide_num, image_path, layout, width=None, height=None):
        """
        Add a slide with a picture to the presentation.
        - slide_num: From 0 start, state which slide you want to add picture
        - image_path: Path to the picture file
        - layout: Layout position for the picture ('right lower corner', 'mid lower', 'left lower corner')
        - width: Optional width of the picture (in centimeters)
        - height: Optional height of the picture (in centimeters)
        """

        if width is None:
            width = self.slide_width / 2  # Default width
        else:
            width = Cm(width)

        if height is None:
            height = self.slide_height / 2  # Default height
        else:
            height = Cm(height)

        # Calculate position based on layout
        if layout == 'right lower corner':
            left = self.slide_width - width
            top = self.slide_height - height
        elif layout == 'mid lower':
            left = (self.slide_width - width) / 2
            top = self.slide_height - height
        elif layout == 'left lower corner':
            left = Cm(0)
            top = self.slide_height - height
        else:
            raise ValueError("Invalid layout option")

        self.prs.slides[slide_num-1].shapes.add_picture(image_path, left, top, width, height)

  def save(self, name):
    self.prs.save(name+".pptx")
    print("File Saved!")
